import sys
import subprocess

# Ensure python-pptx is installed
try:
    import pptx
except ImportError:
    print("Installing python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    import pptx

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

def create_deck():
    prs = Presentation()
    # Set slide dimensions to 16:9 widescreen
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_slide_layout = prs.slide_layouts[6]

    # Color Palette
    PRIMARY_NAVY = RGBColor(10, 25, 47)      # #0A192F
    ACCENT_BLUE = RGBColor(37, 99, 235)      # #2563EB
    EMERALD_GREEN = RGBColor(16, 185, 129)   # #10B981
    AMBER_GOLD = RGBColor(245, 158, 11)     # #F59E0B
    DARK_BG = RGBColor(248, 250, 252)        # #F8FAFC
    CARD_BG = RGBColor(255, 255, 255)
    TEXT_DARK = RGBColor(15, 23, 42)
    TEXT_MUTED = RGBColor(71, 85, 105)

    def add_header(slide, title_text, category_text="SMART INDIA HACKATHON 2026"):
        # Header bar background
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.9))
        tf = header_box.text_frame
        tf.word_wrap = True
        
        p0 = tf.paragraphs[0]
        p0.text = category_text.upper()
        p0.font.size = Pt(10)
        p0.font.bold = True
        p0.font.color.rgb = ACCENT_BLUE

        p1 = tf.add_paragraph()
        p1.text = title_text
        p1.font.size = Pt(22)
        p1.font.bold = True
        p1.font.color.rgb = PRIMARY_NAVY

    # =========================================================================
    # SLIDE 1: TITLE PAGE
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_slide_layout)
    
    # Left Hero Card
    hero_rect = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(0.8), Inches(11.733), Inches(5.9))
    hero_rect.fill.solid()
    hero_rect.fill.fore_color.rgb = PRIMARY_NAVY
    hero_rect.line.color.rgb = PRIMARY_NAVY

    tf1 = hero_rect.text_frame
    tf1.word_wrap = True
    tf1.margin_left = Inches(0.5)
    tf1.margin_top = Inches(0.5)

    p = tf1.paragraphs[0]
    p.text = "SMART INDIA HACKATHON 2026"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = AMBER_GOLD

    p = tf1.add_paragraph()
    p.text = "JEEVAN SETU (जीवन सेतु)"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    p = tf1.add_paragraph()
    p.text = "AI-Based Smart Logistics & Accessibility Intelligence Platform for North Eastern Region (NER)"
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(226, 232, 240)

    p = tf1.add_paragraph()
    p.text = "\nProblem Statement ID: 26002 | Category: Software Edition"
    p.font.size = Pt(14)
    p.font.color.rgb = EMERALD_GREEN
    p.font.bold = True

    p = tf1.add_paragraph()
    p.text = "Theme: Disaster Management & Smart Logistics Infrastructure"
    p.font.size = Pt(13)
    p.font.color.rgb = RGBColor(203, 213, 225)

    p = tf1.add_paragraph()
    p.text = "\nTeam: AlgoForge | Live Prototype: jsalgoforge.netlify.app"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = AMBER_GOLD

    # =========================================================================
    # SLIDE 2: PROPOSED SOLUTION
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide2, "PROPOSED SOLUTION — JEEVAN SETU COMMAND SYSTEM")

    points = [
        ("🏛️ 8-State Unified Command Center", "Connects SDRF, BRO, MDoNER & dispatchers across Assam, Meghalaya, Arunachal Pradesh, Sikkim, Manipur, Mizoram, Nagaland & Tripura."),
        ("🛰️ Sovereign Geospatial Integration", "Real-time road accessibility overlaid with ISRO Bhuvan NRSC Indian Satellite and NASA Earthdata GIBS imagery."),
        ("📵 Offline-First PWA Architecture", "IndexedDB + Service Worker v10.0 enables zero-internet field reporting in Himalayan dead zones (Sela Pass, Haflong)."),
        ("🚑 Smart Priority Supply Dispatch", "Computes alternate routes during cloudbursts/landslides, prioritizing cold-chain vaccines, insulin & fuel over standard cargo."),
        ("🤖 Dual AI Engine Integration", "Google Gemini AI Copilot for incident triage + Scikit-Learn Geotechnical Risk Model (Landslide Hazard & Flood Vulnerability Index)."),
        ("♿ Universal WCAG 2.1 AA Accessibility", "100% WCAG AA compliant with keyboard shortcuts, high contrast, font scaling, and screen-reader announcer.")
    ]

    for i, (title, desc) in enumerate(points):
        row = i // 2
        col = i % 2
        left = Inches(0.8 + col * 5.95)
        top = Inches(1.5 + row * 1.8)
        
        card = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(1.65))
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_BG
        card.line.color.rgb = RGBColor(226, 232, 240)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.15)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_NAVY

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 3: TECHNICAL APPROACH & ARCHITECTURE
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide3, "TECHNICAL APPROACH & DATA FLOW ARCHITECTURE")

    # Architecture Box (Flowchart)
    arch_box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.5), Inches(11.733), Inches(2.2))
    arch_box.fill.solid()
    arch_box.fill.fore_color.rgb = PRIMARY_NAVY
    arch_box.line.color.rgb = PRIMARY_NAVY

    tf_arch = arch_box.text_frame
    tf_arch.word_wrap = True
    tf_arch.margin_left = Inches(0.3)
    tf_arch.margin_top = Inches(0.2)

    p = tf_arch.paragraphs[0]
    p.text = "SYSTEM ARCHITECTURE & END-TO-END DATA FLOW"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = AMBER_GOLD

    p2 = tf_arch.add_paragraph()
    p2.text = "[Field Responders / Mobile PWA]\n  └─► Offline Queue (IndexedDB + Service Worker v10.0)\n        └─► Auto-Sync on Network Reconnect\n              └─► Dual AI Engine (Gemini AI Triage + Scikit-Learn LHI/FVI)\n                    └─► Open-Meteo & ISRO Bhuvan Satellite APIs\n                          └─► OSRM Dynamic Rerouting Engine ──► [MDoNER 8-State Command Radar]"
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(241, 245, 249)

    # Tech Specs Grid below
    specs = [
        ("Frontend & PWA", "HTML5, Tailwind CSS Engine, Vanilla JS, Service Worker v10.0, IndexedDB"),
        ("GIS & Mapping", "Leaflet.js, OpenStreetMap, ISRO Bhuvan NRSC WMS, NASA Earthdata GIBS, Esri Satellite"),
        ("Live APIs & AI", "Open-Meteo Weather API, Nominatim Geocoding API, OSRM Driving Directions API, Google Gemini AI"),
        ("Working Prototype", "Live Deployed Application: jsalgoforge.netlify.app | GitHub: github.com/pragyakatiyar73-cyber/Jeevan-setu")
    ]

    for i, (cat, val) in enumerate(specs):
        left = Inches(0.8 + (i % 2) * 5.95)
        top = Inches(4.0 + (i // 2) * 1.5)
        
        box = slide3.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(1.35))
        box.fill.solid()
        box.fill.fore_color.rgb = DARK_BG
        box.line.color.rgb = RGBColor(226, 232, 240)

        tf = box.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.2)
        tf.margin_top = Inches(0.15)

        p = tf.paragraphs[0]
        p.text = f"🔹 {cat}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_NAVY

        p2 = tf.add_paragraph()
        p2.text = val
        p2.font.size = Pt(10.5)
        p2.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 4: FEASIBILITY AND VIABILITY
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide4, "FEASIBILITY, VIABILITY & RISK MITIGATION")

    feas_cards = [
        ("⚡ Zero Infrastructure Overhead", "Ultra-fast PWA runs on smartphones, tablets, or desktops with zero native installation requirements.", EMERALD_GREEN),
        ("💰 ₹0 Software Subscription Cost", "Keyless, open API architecture (Open-Meteo, Nominatim, OSRM, ISRO Bhuvan) saves millions in government software licensing fees.", ACCENT_BLUE),
        ("📵 Cellular Outage Mitigation", "Total network outage in mountain passes (Sela Pass, Haflong) handled via Service Worker + IndexedDB local offline storage.", AMBER_GOLD),
        ("♿ Inclusivity & Accessibility", "WCAG 2.1 AA compliance ensures usability for field responders with varying digital literacy and visual impairments.", PRIMARY_NAVY)
    ]

    for i, (title, desc, color) in enumerate(feas_cards):
        left = Inches(0.8 + (i % 2) * 5.95)
        top = Inches(1.5 + (i // 2) * 2.6)

        card = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_BG
        card.line.color.rgb = color
        card.line.width = Pt(2)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = color

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11.5)
        p2.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 5: IMPACT AND BENEFITS
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide5, "POTENTIAL IMPACT & MEASURABLE BENEFITS")

    impacts = [
        ("🎯 Target Audience & Reach", "• 1-click regional logistics visibility for MDoNER, NEC, SDRF/NDMA & Border Roads Organisation (BRO).\n• Protects 45+ Million citizens across 8 North Eastern states from severe supply shortages during disaster isolations.", ACCENT_BLUE),
        ("🩺 Social Benefit: Zero Interruption", "• Zero mortality from vaccine, insulin, blood plasma, or medical supply interruptions during seasonal landslides & cloudbursts.", EMERALD_GREEN),
        ("📉 Economic Savings: 35.4% Faster", "• 35.4% reduction in transit delays; massive fuel savings, reduced vehicle breakdowns, and eliminated food spoilage.", AMBER_GOLD),
        ("🌱 Environmental Impact: -24% CO2", "• AI dynamic rerouting avoids blocked mountain sectors, reducing unnecessary detour carbon emissions by 24%.", PRIMARY_NAVY)
    ]

    for i, (title, desc, color) in enumerate(impacts):
        left = Inches(0.8 + (i % 2) * 5.95)
        top = Inches(1.5 + (i // 2) * 2.6)

        card = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, Inches(5.75), Inches(2.4))
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_BG
        card.line.color.rgb = RGBColor(226, 232, 240)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.2)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13.5)
        p.font.bold = True
        p.font.color.rgb = color

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED

    # =========================================================================
    # SLIDE 6: RESEARCH AND REFERENCES
    # =========================================================================
    slide6 = prs.slides.add_slide(blank_slide_layout)
    add_header(slide6, "RESEARCH, REFERENCES & ARTIFACTS")

    refs = [
        ("🏛️ Government & Nodal Ministry Publications", "• MDoNER — Annual Report on NER Infrastructure & Connectivity (mdoner.gov.in)\n• NDMA — Guidelines for Management of Landslides in Himalayan Corridors (ndma.gov.in)\n• BRO — Infrastructure Project Vartak & Beacon Logs (bro.gov.in)"),
        ("🛰️ Geospatial Data Sources & Satellite APIs", "• ISRO Bhuvan NRSC Web Map Service (bhuvan.nrsc.gov.in)\n• NASA GIBS — MODIS & VIIRS Satellite Cloud/Rain Data (earthdata.nasa.gov)\n• Open-Meteo Weather API (open-meteo.com)"),
        ("🚀 Live Project Artifacts & Standards", "• Live Working Prototype: https://jsalgoforge.netlify.app\n• Official GitHub Repository: https://github.com/pragyakatiyar73-cyber/Jeevan-setu\n• Accessibility Compliance: W3C WCAG 2.1 AA Specifications (w3.org/TR/WCAG21)")
    ]

    for i, (title, desc) in enumerate(refs):
        top = Inches(1.5 + i * 1.8)
        card = slide6.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.733), Inches(1.6))
        card.fill.solid()
        card.fill.fore_color.rgb = DARK_BG
        card.line.color.rgb = RGBColor(226, 232, 240)

        tf = card.text_frame
        tf.word_wrap = True
        tf.margin_left = Inches(0.25)
        tf.margin_top = Inches(0.15)

        p = tf.paragraphs[0]
        p.text = title
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_NAVY

        p2 = tf.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(11)
        p2.font.color.rgb = TEXT_MUTED

    output_path = "SIH2026_Jeevan_Setu_AlgoForge.pptx"
    prs.save(output_path)
    print(f"Successfully generated presentation: {output_path}")

if __name__ == "__main__":
    create_deck()
